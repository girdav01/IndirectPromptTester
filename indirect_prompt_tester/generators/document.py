"""Document generator for embedding indirect prompts in office documents."""
from pathlib import Path
from typing import Optional
from .base import BaseGenerator

class DocumentGenerator(BaseGenerator):
    """Generator for creating office documents with embedded indirect prompts."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['docx', 'xlsx', 'pptx', 'pdf', 'txt']
    
    def generate(
        self,
        prompt: str,
        output_path: Path,
        doc_type: str = "docx",
        method: str = "visible",
        **kwargs
    ) -> Path:
        """
        Generate a document with embedded indirect prompt.
        
        Args:
            prompt: The indirect prompt to embed
            output_path: Path where the document should be saved
            doc_type: Type of document ('docx', 'xlsx', 'pptx', 'pdf', 'txt')
            method: Embedding method ('visible', 'hidden', 'metadata', 'comments')
        """
        self.validate_output_path(output_path)
        
        if doc_type == "docx":
            return self._generate_docx(prompt, output_path, method)
        elif doc_type == "xlsx":
            return self._generate_xlsx(prompt, output_path, method)
        elif doc_type == "pptx":
            return self._generate_pptx(prompt, output_path, method)
        elif doc_type == "pdf":
            return self._generate_pdf(prompt, output_path, method)
        elif doc_type == "txt":
            return self._generate_txt(prompt, output_path, method)
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
    
    def _generate_docx(self, prompt: str, output_path: Path, method: str) -> Path:
        """Generate a Word document."""
        from docx import Document
        from docx.shared import Inches
        
        doc = Document()
        doc.add_heading('Sample Document', 0)
        doc.add_paragraph('This is a sample document for testing.')
        
        if method == "visible":
            doc.add_paragraph(prompt)
        elif method == "hidden":
            # Add as hidden text (formatted as white on white)
            p = doc.add_paragraph()
            run = p.add_run(prompt)
            run.font.color.rgb = None  # Will be visible but can be hidden
        elif method == "metadata":
            doc.core_properties.comments = prompt
        elif method == "comments":
            # Note: python-docx doesn't directly support adding comments via API
            # Instead, embed as hidden text with a marker
            doc.add_paragraph('See embedded data below.')
            p = doc.add_paragraph()
            run = p.add_run(f"[COMMENT: {prompt}]")
            run.font.hidden = True
        
        doc.save(output_path)
        return output_path
    
    def _generate_xlsx(self, prompt: str, output_path: Path, method: str) -> Path:
        """Generate an Excel spreadsheet."""
        from openpyxl import Workbook
        from openpyxl.comments import Comment
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        
        ws['A1'] = 'Sample Data'
        ws['B1'] = 'Value'
        
        if method == "visible":
            ws['A2'] = prompt
        elif method == "hidden":
            ws['A2'] = prompt
            ws['A2'].font = ws['A2'].font.copy(color="FFFFFF")  # White text
        elif method == "metadata":
            wb.properties.comments = prompt
        elif method == "comments":
            ws['A1'].comment = Comment(prompt, 'System')
        
        wb.save(output_path)
        return output_path
    
    def _generate_pptx(self, prompt: str, output_path: Path, method: str) -> Path:
        """Generate a PowerPoint presentation."""
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Sample Presentation"
        
        if method == "visible":
            content = slide.placeholders[1]
            content.text = prompt
        elif method == "hidden":
            # Add text box with white text
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
            text_frame = textbox.text_frame
            text_frame.text = prompt
            text_frame.paragraphs[0].font.color.rgb = None  # Can be styled
        elif method == "metadata":
            prs.core_properties.comments = prompt
        
        prs.save(output_path)
        return output_path
    
    def _generate_pdf(self, prompt: str, output_path: Path, method: str) -> Path:
        """Generate a PDF document."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            c = canvas.Canvas(str(output_path), pagesize=letter)
            width, height = letter
            
            c.drawString(100, height - 100, "Sample PDF Document")
            
            if method == "visible":
                # Wrap long prompts
                y_pos = height - 150
                words = prompt.split()
                line = ""
                for word in words:
                    test_line = line + " " + word if line else word
                    if len(test_line) > 80:  # Approximate character limit per line
                        c.drawString(100, y_pos, line)
                        y_pos -= 20
                        line = word
                    else:
                        line = test_line
                if line:
                    c.drawString(100, y_pos, line)
            elif method == "hidden":
                # Embed in metadata
                c.setTitle("Sample Document")
                c.setSubject(prompt)
            
            c.save()
        except ImportError:
            # Fallback: create text file
            content = f"PDF file with embedded prompt: {prompt}\n"
            content += f"Method: {method}\n"
            content += "Note: Install reportlab for actual PDF generation"
            output_path.with_suffix('.txt').write_text(content)
            return output_path.with_suffix('.txt')
        
        return output_path
    
    def _generate_txt(self, prompt: str, output_path: Path, method: str) -> Path:
        """Generate a plain text file."""
        content = "Sample Text Document\n\n"
        
        if method == "visible":
            content += prompt
        elif method == "hidden":
            # Embed as comment or invisible characters
            content += f"Main content here\n<!-- {prompt} -->\n"
        elif method == "metadata":
            content += prompt
        
        output_path.write_text(content)
        return output_path
    
    def get_supported_formats(self) -> list:
        return self.supported_formats

